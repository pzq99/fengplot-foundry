#!/usr/bin/env python3
"""Convert Matplotlib SVG panels to native editable PowerPoint shapes and text."""

from __future__ import annotations

import argparse
import json
import math
import re
import shutil
import tempfile
from collections import Counter
from dataclasses import dataclass
from pathlib import Path

from lxml import etree
from svgpathtools import Arc, CubicBezier, Line, QuadraticBezier, parse_path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Inches, Pt


SVG_NS = "http://www.w3.org/2000/svg"
XLINK_NS = "http://www.w3.org/1999/xlink"
EMU_PER_INCH = 914400
PT_PER_INCH = 72.0
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


Matrix = tuple[float, float, float, float, float, float]
IDENTITY: Matrix = (1, 0, 0, 1, 0, 0)


@dataclass(frozen=True)
class PanelGeometry:
    path: Path
    view_box: tuple[float, float, float, float]

    @property
    def aspect(self) -> float:
        return self.view_box[2] / self.view_box[3]


def local_name(element) -> str:
    return element.tag.rsplit("}", 1)[-1] if isinstance(element.tag, str) else ""


def multiply(m1: Matrix, m2: Matrix) -> Matrix:
    a1, b1, c1, d1, e1, f1 = m1
    a2, b2, c2, d2, e2, f2 = m2
    return (
        a1 * a2 + c1 * b2,
        b1 * a2 + d1 * b2,
        a1 * c2 + c1 * d2,
        b1 * c2 + d1 * d2,
        a1 * e2 + c1 * f2 + e1,
        b1 * e2 + d1 * f2 + f1,
    )


def transform_point(m: Matrix, x: float, y: float) -> tuple[float, float]:
    a, b, c, d, e, f = m
    return a * x + c * y + e, b * x + d * y + f


def parse_numbers(value: str) -> list[float]:
    return [float(x) for x in re.findall(r"[-+]?(?:\d*\.\d+|\d+\.?)(?:[eE][-+]?\d+)?", value)]


def clip_polygon_to_rect(
    points: list[tuple[float, float]], rect: tuple[float, float, float, float]
) -> list[tuple[float, float]]:
    """Clip a closed polygon to an axis-aligned SVG clip rectangle."""
    x_min, y_min, x_max, y_max = rect

    def clip_edge(vertices, inside, intersect):
        if not vertices:
            return []
        output = []
        previous = vertices[-1]
        previous_inside = inside(previous)
        for current in vertices:
            current_inside = inside(current)
            if current_inside:
                if not previous_inside:
                    output.append(intersect(previous, current))
                output.append(current)
            elif previous_inside:
                output.append(intersect(previous, current))
            previous, previous_inside = current, current_inside
        return output

    def at_x(a, b, x):
        if abs(b[0] - a[0]) < 1e-12:
            return x, a[1]
        t = (x - a[0]) / (b[0] - a[0])
        return x, a[1] + t * (b[1] - a[1])

    def at_y(a, b, y):
        if abs(b[1] - a[1]) < 1e-12:
            return a[0], y
        t = (y - a[1]) / (b[1] - a[1])
        return a[0] + t * (b[0] - a[0]), y

    clipped = points
    clipped = clip_edge(clipped, lambda p: p[0] >= x_min, lambda a, b: at_x(a, b, x_min))
    clipped = clip_edge(clipped, lambda p: p[0] <= x_max, lambda a, b: at_x(a, b, x_max))
    clipped = clip_edge(clipped, lambda p: p[1] >= y_min, lambda a, b: at_y(a, b, y_min))
    clipped = clip_edge(clipped, lambda p: p[1] <= y_max, lambda a, b: at_y(a, b, y_max))
    return clipped


def clip_polyline_to_rect(
    points: list[tuple[float, float]], rect: tuple[float, float, float, float]
) -> list[list[tuple[float, float]]]:
    """Clip an open polyline to an axis-aligned SVG clip rectangle.

    A clipped line can leave and re-enter the rectangle, so the result is a
    list of independent runs rather than a single point list.
    """
    if len(points) < 2:
        return []
    x_min, y_min, x_max, y_max = rect

    def clip_segment(a, b):
        dx, dy = b[0] - a[0], b[1] - a[1]
        lower, upper = 0.0, 1.0
        for direction, distance in (
            (-dx, a[0] - x_min),
            (dx, x_max - a[0]),
            (-dy, a[1] - y_min),
            (dy, y_max - a[1]),
        ):
            if abs(direction) < 1e-12:
                if distance < 0:
                    return None
                continue
            ratio = distance / direction
            if direction < 0:
                lower = max(lower, ratio)
            else:
                upper = min(upper, ratio)
            if lower > upper:
                return None
        return (
            (a[0] + lower * dx, a[1] + lower * dy),
            (a[0] + upper * dx, a[1] + upper * dy),
        )

    runs: list[list[tuple[float, float]]] = []
    current: list[tuple[float, float]] = []
    for start, end in zip(points, points[1:]):
        clipped = clip_segment(start, end)
        if clipped is None:
            if len(current) >= 2:
                runs.append(current)
            current = []
            continue
        clipped_start, clipped_end = clipped
        if current and math.dist(current[-1], clipped_start) <= 1e-7:
            if math.dist(current[-1], clipped_end) > 1e-7:
                current.append(clipped_end)
        else:
            if len(current) >= 2:
                runs.append(current)
            current = [clipped_start, clipped_end]
    if len(current) >= 2:
        runs.append(current)
    return runs


def parse_transform(value: str | None) -> Matrix:
    if not value:
        return IDENTITY
    result = IDENTITY
    for name, args_text in re.findall(r"([A-Za-z]+)\s*\(([^)]*)\)", value):
        args = parse_numbers(args_text)
        name = name.lower()
        if name == "matrix" and len(args) == 6:
            op = tuple(args)  # type: ignore[assignment]
        elif name == "translate":
            op = (1, 0, 0, 1, args[0], args[1] if len(args) > 1 else 0)
        elif name == "scale":
            sy = args[1] if len(args) > 1 else args[0]
            op = (args[0], 0, 0, sy, 0, 0)
        elif name == "rotate":
            angle = math.radians(args[0])
            rot = (math.cos(angle), math.sin(angle), -math.sin(angle), math.cos(angle), 0, 0)
            if len(args) >= 3:
                cx, cy = args[1], args[2]
                op = multiply(multiply((1, 0, 0, 1, cx, cy), rot), (1, 0, 0, 1, -cx, -cy))
            else:
                op = rot
        elif name == "skewx":
            op = (1, 0, math.tan(math.radians(args[0])), 1, 0, 0)
        elif name == "skewy":
            op = (1, math.tan(math.radians(args[0])), 0, 1, 0, 0)
        else:
            continue
        result = multiply(result, op)
    return result


def parse_style(element) -> dict[str, str]:
    result: dict[str, str] = {}
    raw = element.get("style") or ""
    for item in raw.split(";"):
        if ":" in item:
            key, value = item.split(":", 1)
            result[key.strip()] = value.strip()
    for key in (
        "fill",
        "fill-opacity",
        "stroke",
        "stroke-opacity",
        "stroke-width",
        "stroke-dasharray",
        "opacity",
        "font-size",
        "font-style",
        "font-weight",
        "font-family",
        "text-anchor",
    ):
        if element.get(key) is not None:
            result[key] = element.get(key)
    return result


NAMED_COLORS = {
    "black": "000000",
    "white": "FFFFFF",
    "red": "FF0000",
    "green": "008000",
    "blue": "0000FF",
    "gray": "808080",
    "grey": "808080",
    "orange": "FFA500",
    "purple": "800080",
}


def color_value(value: str | None) -> RGBColor | None:
    if not value or value.lower() in {"none", "transparent"}:
        return None
    value = value.strip().lower()
    if value.startswith("#"):
        raw = value[1:]
        if len(raw) == 3:
            raw = "".join(x * 2 for x in raw)
        if len(raw) == 6:
            return RGBColor.from_string(raw.upper())
    if value.startswith("rgb"):
        nums = parse_numbers(value)
        if len(nums) >= 3:
            return RGBColor(*(max(0, min(255, round(x))) for x in nums[:3]))
    if value in NAMED_COLORS:
        return RGBColor.from_string(NAMED_COLORS[value])
    return RGBColor(0, 0, 0)


def numeric(value: str | None, default: float = 0.0) -> float:
    if not value:
        return default
    numbers = parse_numbers(value)
    return numbers[0] if numbers else default


def opacity_value(value: str | None, default: float = 1.0) -> float:
    """Return an SVG opacity as a clamped 0..1 value."""
    if not value:
        return default
    result = numeric(value, default)
    if "%" in value:
        result /= 100.0
    return max(0.0, min(1.0, result))


def set_color_opacity(color_format, opacity: float) -> None:
    """Write SVG alpha to the DrawingML color backing a python-pptx color."""
    opacity = max(0.0, min(1.0, opacity))
    if opacity >= 0.9999:
        return
    color_element = color_format._color._xClr
    for child in list(color_element):
        if local_name(child) == "alpha":
            color_element.remove(child)
    alpha = OxmlElement("a:alpha")
    # DrawingML alpha is opacity (not transparency), expressed in 1/1000 percent.
    alpha.set("val", str(round(opacity * 100000)))
    color_element.append(alpha)


def read_panel_geometry(path: Path) -> PanelGeometry:
    parser = etree.XMLParser(resolve_entities=False, no_network=True, recover=False)
    root = etree.parse(str(path), parser).getroot()
    values = parse_numbers(root.get("viewBox", ""))
    if len(values) != 4:
        values = [0.0, 0.0, numeric(root.get("width"), 1.0), numeric(root.get("height"), 1.0)]
    if values[2] <= 0 or values[3] <= 0:
        raise ValueError(f"Invalid SVG canvas in {path}: {values}")
    return PanelGeometry(path, tuple(values))  # type: ignore[arg-type]


def read_panel_title(path: Path) -> str:
    parser = etree.XMLParser(resolve_entities=False, no_network=True, recover=False)
    root = etree.parse(str(path), parser).getroot()
    for candidate in root.iter():
        if local_name(candidate) != "text":
            continue
        text = "".join(candidate.itertext()).strip()
        if is_panel_title(text):
            return text
    return path.stem


def balanced_row_counts(panel_count: int, max_columns: int) -> list[int]:
    """Return stable row counts such as 7 -> 3/2/2 and 8 -> 3/3/2."""
    row_count = math.ceil(panel_count / max_columns)
    base, extra = divmod(panel_count, row_count)
    return [base + (1 if index < extra else 0) for index in range(row_count)]


def adaptive_layout(
    panels: list[PanelGeometry],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    max_columns: int,
    min_gap_x: float,
    min_gap_y: float,
    legend_band_height: float = 0.0,
    lock_common_panel_size: bool = False,
    stretch_row_panels: bool = False,
    stretch_rows: set[int] | None = None,
    stretch_row_weights: list[float] | None = None,
    row_counts: list[int] | None = None,
) -> list[tuple[float, float, float, float]]:
    """Fit heterogeneous panel canvases without stretching their aspect ratios.

    Rows are balanced, row tops/bottoms are shared, the first and last panel in
    every multi-panel row touch the global horizontal boundaries, and remaining
    whitespace becomes equal inter-panel gaps.
    """
    if row_counts is None:
        counts = balanced_row_counts(len(panels), max_columns)
    else:
        counts = list(row_counts)
        if not counts or any(count < 1 or count > max_columns for count in counts):
            raise ValueError(
                f"Every explicit row count must be between 1 and {max_columns}: {counts}"
            )
        if sum(counts) != len(panels):
            raise ValueError(
                f"Explicit row counts sum to {sum(counts)}, but {len(panels)} SVG panels were found"
            )
    rows: list[list[PanelGeometry]] = []
    cursor = 0
    for count in counts:
        rows.append(panels[cursor : cursor + count])
        cursor += count

    usable_height = (
        height
        - min_gap_y * max(0, len(rows) - 1)
        - (legend_band_height if len(rows) > 1 else 0.0)
    )
    if usable_height <= 0:
        raise ValueError("Vertical gaps leave no room for SVG panels")
    band_height = usable_height / len(rows)
    row_top = top

    boxes: list[tuple[float, float, float, float]] = []
    aspect_counts = Counter(round(panel.aspect, 6) for panel in panels)
    common_aspect = aspect_counts.most_common(1)[0][0]
    common_height = min(
        band_height,
        (width - min_gap_x * max(0, max_columns - 1)) / (max_columns * common_aspect),
    )

    for row_index, row in enumerate(rows):
        usable_width = width - min_gap_x * max(0, len(row) - 1)
        if usable_width <= 0:
            raise ValueError("Horizontal gaps leave no room for SVG panels")
        if stretch_rows and row_index + 1 in stretch_rows:
            # Fill a selected row horizontally while keeping its full vertical
            # band height. Widths remain proportional to source aspect ratios,
            # so a deliberately wider panel (for example a heatmap) stays the
            # widest item in the row.
            if stretch_row_weights is not None and len(stretch_row_weights) != len(row):
                raise ValueError(
                    f"--stretch-row-weights has {len(stretch_row_weights)} values, "
                    f"but selected row {row_index + 1} has {len(row)} panels"
                )
            weights = stretch_row_weights or [1.0 for _ in row]
            weighted_aspects = [panel.aspect * weight for panel, weight in zip(row, weights)]
            aspect_sum = sum(weighted_aspects)
            widths = [usable_width * value / aspect_sum for value in weighted_aspects]
            heights = [band_height for _ in row]
            panel_tops = [row_top for _ in row]
        elif stretch_row_panels:
            # Presentation-first option for dense 3-up figures: every panel gets
            # an equal-width cell so the populated row spans the complete content
            # region. The user explicitly accepts horizontal stretching here.
            widths = [usable_width / len(row) for _ in row]
            heights = [common_height for _ in row]
            panel_tops = [row_top for _ in row]
        elif lock_common_panel_size:
            heights = [common_height for _ in row]
            widths = [panel.aspect * common_height for panel in row]
            overflow = sum(widths) - usable_width
            if overflow > 0:
                exceptional = [
                    index
                    for index, panel in enumerate(row)
                    if abs(panel.aspect - common_aspect) > 1e-4
                ]
                shrinkable = exceptional or list(range(len(row)))
                remaining = overflow
                for index in shrinkable:
                    share = remaining / (len(shrinkable) - shrinkable.index(index))
                    minimum_width = widths[index] * 0.55
                    reduction = min(share, widths[index] - minimum_width)
                    widths[index] -= reduction
                    heights[index] = widths[index] / row[index].aspect
                    remaining -= reduction
                if remaining > 1e-6:
                    scale = usable_width / sum(widths)
                    widths = [value * scale for value in widths]
                    heights = [value * scale for value in heights]
            panel_tops = [row_top for _ in row]
        else:
            fit_height = usable_width / sum(panel.aspect for panel in row)
            row_height = min(band_height, fit_height)
            heights = [row_height for _ in row]
            panel_tops = [row_top + (band_height - row_height) / 2 for _ in row]
            widths = [panel.aspect * row_height for panel in row]
        if len(row) == 1:
            x_positions = [left + (width - widths[0]) / 2]
        else:
            gap_x = (width - sum(widths)) / (len(row) - 1)
            x_positions = []
            x = left
            for panel_width in widths:
                x_positions.append(x)
                x += panel_width + gap_x
        boxes.extend(
            (panel_left, panel_top, panel_width, panel_height)
            for panel_left, panel_top, panel_width, panel_height in zip(
                x_positions, panel_tops, widths, heights
            )
        )
        row_top += band_height + min_gap_y
        if row_index == 0 and len(rows) > 1:
            row_top += legend_band_height
    return boxes


def is_panel_title(text: str) -> bool:
    return bool(re.match(r"^\s*(?:\([A-Za-z]\)|[A-Za-z][.)])\s*\S", text))


def set_shape_name(shape, value: str) -> None:
    c_nv_pr = shape._element.find(f".//{{{P_NS}}}cNvPr")
    if c_nv_pr is not None:
        c_nv_pr.set("name", value)


class SvgPanelRenderer:
    def __init__(
        self,
        slide,
        svg_path: Path,
        box: tuple[float, float, float, float],
        target_font_pt: float,
        font_name: str,
        normalize_titles: bool = True,
        preserve_source_font_ratios: bool = False,
        preserved_ratio_font_levels: tuple[float, float, float] | None = None,
        preserved_ratio_role_levels: tuple[float, float, float, float] | None = None,
        normalize_axis_box: tuple[float, float, float, float] | None = None,
        suppress_legends: bool = False,
        compact_rotated_labels: bool = False,
        title_font_scale: float = 1.0,
        font_levels: tuple[float, float] | None = None,
        dense_font_pt: float | None = None,
        very_dense_font_pt: float | None = None,
        min_stroke_width_pt: float | None = None,
        axis_stroke_width_pt: float | None = None,
        suppressed_texts: set[str] | None = None,
        category_proxy_items: list[dict] | None = None,
        category_proxy_gap_in: float = 0.60,
    ):
        self.slide = slide
        self.svg_path = svg_path
        self.left_in, self.top_in, self.width_in, self.height_in = box
        self.target_font_pt = target_font_pt
        self.font_name = font_name
        self.normalize_titles = normalize_titles
        self.preserve_source_font_ratios = preserve_source_font_ratios
        self.preserved_ratio_font_levels = preserved_ratio_font_levels
        self.preserved_ratio_role_levels = preserved_ratio_role_levels
        self.normalize_axis_box = normalize_axis_box
        self.suppress_legends = suppress_legends
        self.compact_rotated_labels = compact_rotated_labels
        self.title_font_scale = title_font_scale
        self.font_levels = font_levels
        self.dense_font_pt = dense_font_pt
        self.very_dense_font_pt = very_dense_font_pt
        self.min_stroke_width_pt = min_stroke_width_pt
        self.axis_stroke_width_pt = axis_stroke_width_pt
        self.suppressed_texts = suppressed_texts or set()
        self.category_proxy_items = category_proxy_items or []
        self.category_proxy_gap_in = category_proxy_gap_in
        parser = etree.XMLParser(resolve_entities=False, no_network=True, recover=False)
        self.root = etree.parse(str(svg_path), parser).getroot()
        view_box = parse_numbers(self.root.get("viewBox", ""))
        if len(view_box) != 4:
            width = numeric(self.root.get("width"), 1)
            height = numeric(self.root.get("height"), 1)
            view_box = [0, 0, width, height]
        self.vx, self.vy, self.vw, self.vh = view_box
        self.sx_emu = self.width_in * EMU_PER_INCH / self.vw
        self.sy_emu = self.height_in * EMU_PER_INCH / self.vh
        self.ids = {e.get("id"): e for e in self.root.iter() if e.get("id")}
        self.clip_rects: dict[str, tuple[float, float, float, float]] = {}
        for candidate in self.root.iter():
            if local_name(candidate) != "clipPath" or not candidate.get("id"):
                continue
            rect = next((child for child in candidate if local_name(child) == "rect"), None)
            if rect is None:
                continue
            x = numeric(rect.get("x"))
            y = numeric(rect.get("y"))
            width = numeric(rect.get("width"))
            height = numeric(rect.get("height"))
            if width > 0 and height > 0:
                self.clip_rects[candidate.get("id")] = (x, y, x + width, y + height)
        self.axis_normalization_matrix: Matrix | None = None
        if self.normalize_axis_box and self.clip_rects:
            source = max(
                self.clip_rects.values(),
                key=lambda rect: (rect[2] - rect[0]) * (rect[3] - rect[1]),
            )
            left_f, right_f, top_f, bottom_f = self.normalize_axis_box
            target = (
                self.vx + left_f * self.vw,
                self.vy + top_f * self.vh,
                self.vx + right_f * self.vw,
                self.vy + bottom_f * self.vh,
            )
            scale_x = (target[2] - target[0]) / (source[2] - source[0])
            scale_y = (target[3] - target[1]) / (source[3] - source[1])
            self.axis_normalization_matrix = (
                scale_x,
                0,
                0,
                scale_y,
                target[0] - source[0] * scale_x,
                target[1] - source[1] * scale_y,
            )
        source_font_sizes = []
        title_font_sizes = []
        for candidate in self.root.iter():
            if local_name(candidate) != "text":
                continue
            candidate_style = parse_style(candidate)
            candidate_size = numeric(candidate_style.get("font-size"), 0.0)
            if candidate_size > 0:
                source_font_sizes.append(candidate_size)
                if is_panel_title("".join(candidate.itertext()).strip()):
                    title_font_sizes.append(candidate_size)
        self.source_base_font_pt = (
            max(title_font_sizes) if title_font_sizes else max(source_font_sizes, default=34.0)
        )
        self.dense_panel = len(
            self.root.xpath('//*[starts-with(@id, "axes_")]')
        ) > 1
        self.very_dense_panel = len(self.root.xpath('//*[local-name()="text"]')) >= 35
        self.created = []
        self.image_count = 0

    def add_panel_anchor(self) -> None:
        anchor = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(self.left_in),
            Inches(self.top_in),
            Inches(self.width_in),
            Inches(self.height_in),
        )
        anchor.fill.background()
        anchor.line.fill.background()
        set_shape_name(anchor, f"Panel layout anchor: {self.svg_path.stem}")
        self.created.append(anchor)

    def add_common_axis_anchor(self) -> None:
        if not self.normalize_axis_box:
            return
        left_f, right_f, top_f, bottom_f = self.normalize_axis_box
        left, top = self.map_point((self.vx + left_f * self.vw, self.vy + top_f * self.vh))
        right, bottom = self.map_point((self.vx + right_f * self.vw, self.vy + bottom_f * self.vh))
        anchor = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left,
            top,
            right - left,
            bottom - top,
        )
        anchor.fill.background()
        anchor.line.fill.background()
        set_shape_name(anchor, f"Common axes anchor: {self.svg_path.stem}")
        self.created.append(anchor)

    def add_category_proxy_markers(self) -> None:
        """Add color-only column keys above a categorical heatmap.

        The shared legend already carries method names. These compact markers
        preserve the column-to-method mapping after redundant rotated tick
        labels are suppressed.
        """
        if not self.category_proxy_items or not self.clip_rects:
            return
        grid = max(
            self.clip_rects.values(),
            key=lambda rect: (rect[2] - rect[0]) * (rect[3] - rect[1]),
        )
        count = len(self.category_proxy_items)
        marker_size_in = 0.42
        grid_top = self.map_point((grid[0], grid[1]))[1]
        marker_top = grid_top - Inches(marker_size_in + self.category_proxy_gap_in)
        cell_width = (grid[2] - grid[0]) / count
        for index, item in enumerate(self.category_proxy_items):
            center_svg_x = grid[0] + (index + 0.5) * cell_width
            center_x = self.map_point((center_svg_x, grid[1]))[0]
            marker = self.slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                center_x - Inches(marker_size_in / 2),
                marker_top,
                Inches(marker_size_in),
                Inches(marker_size_in),
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = color_value(str(item["color"])) or RGBColor(0, 0, 0)
            marker.line.color.rgb = RGBColor.from_string(
                str(item.get("line_color", "333333")).lstrip("#")
            )
            marker.line.width = Pt(float(item.get("line_width", 1.5)))
            set_shape_name(
                marker,
                f"Category proxy marker {self.svg_path.stem} {index + 1}: {item['label']}",
            )
            self.created.append(marker)

    def map_point(self, p: tuple[float, float]) -> tuple[int, int]:
        x = (p[0] - self.vx) * self.sx_emu + self.left_in * EMU_PER_INCH
        y = (p[1] - self.vy) * self.sy_emu + self.top_in * EMU_PER_INCH
        return round(x), round(y)

    @property
    def stroke_scale_pt(self) -> float:
        return ((self.sx_emu + self.sy_emu) / 2) / 12700.0

    def apply_style(self, shape, style: dict[str, str], closed: bool) -> None:
        element_opacity = opacity_value(style.get("opacity"))
        fill_key = style.get("fill")
        fill_color = color_value(fill_key)
        if fill_key is None:
            # SVG's initial fill value is black. Matplotlib relies on that
            # default for outlined annotation glyphs, while open paths remain
            # unfilled. Explicit fill:none continues to map to no fill.
            fill_color = RGBColor(0, 0, 0) if closed else None
        # python-pptx Connector objects expose a line but no fill. SVG <line>
        # elements legitimately reach this shared styling path, so skip fill
        # handling for connector-backed shapes while preserving it for paths,
        # rectangles, ellipses, and other closed geometry.
        if hasattr(shape, "fill"):
            if fill_color is None:
                shape.fill.background()
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = fill_color
                set_color_opacity(
                    shape.fill.fore_color,
                    element_opacity * opacity_value(style.get("fill-opacity")),
                )
        stroke_key = style.get("stroke")
        stroke_color = color_value(stroke_key)
        if stroke_key is None:
            stroke_color = None
        if stroke_color is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = stroke_color
            set_color_opacity(
                shape.line.color,
                element_opacity * opacity_value(style.get("stroke-opacity")),
            )
            width = numeric(style.get("stroke-width"), 1.0) * self.stroke_scale_pt
            if self.axis_stroke_width_pt is not None and stroke_color == RGBColor(0, 0, 0):
                width = self.axis_stroke_width_pt
            elif self.min_stroke_width_pt is not None:
                width = max(width, self.min_stroke_width_pt)
            shape.line.width = Pt(max(0.5, width))
            dash = style.get("stroke-dasharray")
            if dash and dash.lower() != "none":
                shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    def add_path(
        self,
        d: str,
        matrix: Matrix,
        style: dict[str, str],
        clip_rect: tuple[float, float, float, float] | None = None,
    ) -> None:
        try:
            path = parse_path(d)
        except Exception as exc:
            raise ValueError(f"Cannot parse SVG path in {self.svg_path}: {exc}") from exc
        subpaths = path.continuous_subpaths()
        if not subpaths:
            return
        contours: list[tuple[list[tuple[float, float]], bool]] = []
        for subpath in subpaths:
            if not subpath:
                continue
            points = [transform_point(matrix, subpath[0].start.real, subpath[0].start.imag)]
            for segment in subpath:
                steps = 1 if isinstance(segment, Line) else 8
                if isinstance(segment, (CubicBezier, QuadraticBezier, Arc)):
                    try:
                        steps = max(4, min(20, math.ceil(segment.length(error=1e-3) / 10)))
                    except Exception:
                        steps = 8
                for step in range(1, steps + 1):
                    point = segment.point(step / steps)
                    points.append(transform_point(matrix, point.real, point.imag))
            closed = abs(subpath.start - subpath.end) < 1e-6
            if clip_rect is not None:
                if closed:
                    points = clip_polygon_to_rect(points, clip_rect)
                    if len(points) < 3:
                        continue
                    contours.append((points, True))
                else:
                    contours.extend((run, False) for run in clip_polyline_to_rect(points, clip_rect))
            else:
                contours.append((points, closed))
        if not contours:
            return
        first = contours[0][0][0]
        builder = self.slide.shapes.build_freeform(first[0], first[1], scale=(self.sx_emu, self.sy_emu))
        for index, (points, closed) in enumerate(contours):
            if index:
                builder.move_to(points[0][0], points[0][1])
            builder.add_line_segments(points[1:], close=closed)
        shape = builder.convert_to_shape(
            origin_x=Emu(round(self.left_in * EMU_PER_INCH - self.vx * self.sx_emu)),
            origin_y=Emu(round(self.top_in * EMU_PER_INCH - self.vy * self.sy_emu)),
        )
        self.apply_style(shape, style, any(closed for _, closed in contours))
        self.created.append(shape)

    def add_poly(self, points_text: str, matrix: Matrix, style: dict[str, str], closed: bool) -> None:
        values = parse_numbers(points_text)
        points = [transform_point(matrix, values[i], values[i + 1]) for i in range(0, len(values) - 1, 2)]
        if len(points) < 2:
            return
        builder = self.slide.shapes.build_freeform(points[0][0], points[0][1], scale=(self.sx_emu, self.sy_emu))
        builder.add_line_segments(points[1:], close=closed)
        shape = builder.convert_to_shape(
            origin_x=Emu(round(self.left_in * EMU_PER_INCH - self.vx * self.sx_emu)),
            origin_y=Emu(round(self.top_in * EMU_PER_INCH - self.vy * self.sy_emu)),
        )
        self.apply_style(shape, style, closed)
        self.created.append(shape)

    def add_line(self, element, matrix: Matrix, style: dict[str, str]) -> None:
        p1 = transform_point(matrix, numeric(element.get("x1")), numeric(element.get("y1")))
        p2 = transform_point(matrix, numeric(element.get("x2")), numeric(element.get("y2")))
        x1, y1 = self.map_point(p1)
        x2, y2 = self.map_point(p2)
        shape = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        self.apply_style(shape, style, False)
        self.created.append(shape)

    def add_ellipse(self, element, matrix: Matrix, style: dict[str, str], circle: bool) -> None:
        cx = numeric(element.get("cx"))
        cy = numeric(element.get("cy"))
        rx = numeric(element.get("r")) if circle else numeric(element.get("rx"))
        ry = rx if circle else numeric(element.get("ry"))
        points = []
        for index in range(33):
            theta = 2 * math.pi * index / 32
            points.append(transform_point(matrix, cx + rx * math.cos(theta), cy + ry * math.sin(theta)))
        self.add_poly(" ".join(f"{x},{y}" for x, y in points), IDENTITY, style, True)

    def add_rect(self, element, matrix: Matrix, style: dict[str, str]) -> None:
        x, y = numeric(element.get("x")), numeric(element.get("y"))
        w, h = numeric(element.get("width")), numeric(element.get("height"))
        corners = [
            transform_point(matrix, x, y),
            transform_point(matrix, x + w, y),
            transform_point(matrix, x + w, y + h),
            transform_point(matrix, x, y + h),
        ]
        self.add_poly(" ".join(f"{px},{py}" for px, py in corners), IDENTITY, style, True)

    def add_text(self, element, matrix: Matrix, style: dict[str, str]) -> None:
        text = "".join(element.itertext()).strip()
        if not text:
            return
        if text in self.suppressed_texts:
            return
        x = numeric(element.get("x"))
        y = numeric(element.get("y"))
        px, py = transform_point(matrix, x, y)
        panel_title = self.normalize_titles and is_panel_title(text)
        source_font_pt = numeric(style.get("font-size"), self.source_base_font_pt)
        panel_subtitle = (
            not panel_title
            and self.vx <= x <= self.vx + self.vw * 0.12
            and self.vy <= y <= self.vy + self.vh * 0.18
            and self.source_base_font_pt * 0.45 <= source_font_pt <= self.source_base_font_pt * 0.75
            and len(text) >= 12
        )
        if panel_title:
            # Matplotlib normally aligns an axes title to the plot spine. The
            # delivery standard aligns panel headings to the full panel canvas.
            px = self.vx
            style = dict(style)
            style["text-anchor"] = "start"
        anchor_x, baseline_y = self.map_point((px, py))
        if panel_title:
            # Use a true object offset instead of a text-frame margin. PowerPoint
            # can scale group-internal margins inconsistently, while the native
            # 0.40-inch offset survives grouping and keeps 88 pt italic ink off the
            # delivery boundary on every renderer.
            anchor_x += Inches(0.40)
        angle = math.degrees(math.atan2(matrix[1], matrix[0]))
        anchor = style.get("text-anchor", "start")
        numeric_label = bool(re.fullmatch(r"[−–-]?\d+(?:\.\d+)?%?", text))
        categorical_tick = (
            not panel_title
            and not panel_subtitle
            and abs(angle) < 0.5
            and anchor == "end"
            and not numeric_label
        )
        ratio = source_font_pt / self.source_base_font_pt
        font_pt = self.target_font_pt
        if self.font_levels:
            regular_pt, detail_pt = self.font_levels
            if panel_title:
                font_pt = regular_pt
            else:
                # Presentation-first mode: every non-heading chart label uses
                # one common readable size. This prevents mixed micro-fonts
                # across ticks, annotations, legends, heatmaps, and composites.
                font_pt = detail_pt
                if self.very_dense_panel and self.very_dense_font_pt is not None:
                    font_pt = self.very_dense_font_pt
                elif self.dense_panel and self.dense_font_pt is not None:
                    font_pt = self.dense_font_pt
        elif self.preserve_source_font_ratios:
            # Source labels close to the panel-heading size are primary text and
            # stay at the required target size. Smaller ticks, legends, and cell
            # annotations retain their intentional hierarchy and avoid overflow.
            if self.preserved_ratio_role_levels and not panel_title:
                # Standardized upstream SVGs use source sizes at approximately
                # 85%, 75%, 60%, and <50% of the panel-heading size for section
                # headings, axis labels, body/ticks, and micro detail. Quantize
                # those roles to a small intentional PowerPoint type system.
                section_pt, axis_pt, body_pt, micro_pt = (
                    self.preserved_ratio_role_levels
                )
                if ratio >= 0.80:
                    font_pt = section_pt
                elif ratio >= 0.68:
                    font_pt = axis_pt
                elif ratio >= 0.50:
                    font_pt = body_pt
                else:
                    font_pt = micro_pt
            elif self.preserved_ratio_font_levels and not panel_title:
                primary_pt, secondary_pt, micro_pt = (
                    self.preserved_ratio_font_levels
                )
                if ratio >= 0.58:
                    font_pt = primary_pt
                elif ratio >= 0.33:
                    font_pt = secondary_pt
                else:
                    font_pt = micro_pt
            elif ratio < 0.9:
                font_pt = self.target_font_pt * ratio
        if panel_title:
            font_pt *= self.title_font_scale
        if (
            panel_subtitle
            and not self.font_levels
            and not self.preserved_ratio_role_levels
        ):
            available_width_pt = max(
                font_pt * 3,
                (Inches(self.left_in + self.width_in - 0.20) - anchor_x) / 12700,
            )
            estimated_width_pt = font_pt * (len(text) * 0.57 + 0.25)
            if estimated_width_pt > available_width_pt:
                font_pt *= available_width_pt / estimated_width_pt
        if (
            not self.font_levels
            and not self.preserved_ratio_role_levels
            and self.compact_rotated_labels
            and not panel_title
            and 5 < abs(angle) < 85
        ):
            # Exception panels such as narrow heatmaps often carry several long
            # category labels. A strong compacting factor keeps adjacent rotated
            # labels distinct after the large-format up-scale.
            font_pt *= 0.58
        char_width_pt = font_pt * 0.57
        width_pt = max(font_pt * 1.9, len(text) * char_width_pt + font_pt * 0.25)
        if panel_subtitle and anchor == "start":
            # Hand-tuned reference decks keep long subtitle objects inside the
            # panel boundary without shrinking the common plot font. PowerPoint
            # renders the unwrapped text correctly while the bounded text box
            # prevents group bounds from leaking outside the delivery canvas.
            available_width_pt = max(
                font_pt * 1.9,
                (Inches(self.left_in + self.width_in - 0.02) - anchor_x) / 12700,
            )
            width_pt = min(width_pt, available_width_pt)
        height_pt = font_pt * 1.28
        if abs(abs(angle) - 90) < 0.5:
            # Use native vertical text rather than a rotated wide box. This keeps the
            # PowerPoint object's own bounds inside the allotted panel and avoids the
            # negative-x text boxes found in many manually converted SVG decks.
            box_width_pt = height_pt
            box_height_pt = width_pt
            left = anchor_x - round(box_width_pt * 12700 / 2)
            top = baseline_y - round(box_height_pt * 12700 / 2)
            shape = self.slide.shapes.add_textbox(
                left, top, round(box_width_pt * 12700), round(box_height_pt * 12700)
            )
            alignment = PP_ALIGN.CENTER
            vertical_text = True
        else:
            vertical_text = False
            width_emu = round(width_pt * 12700)
            height_emu = round(height_pt * 12700)
            baseline_offset = font_pt * 0.34 * 12700
            if anchor == "middle":
                anchor_vector_x = 0
                alignment = PP_ALIGN.CENTER
            elif anchor == "end":
                anchor_vector_x = width_emu / 2
                alignment = PP_ALIGN.RIGHT
            else:
                anchor_vector_x = -width_emu / 2
                alignment = PP_ALIGN.LEFT
            theta = math.radians(angle)
            rotated_x = math.cos(theta) * anchor_vector_x - math.sin(theta) * baseline_offset
            rotated_y = math.sin(theta) * anchor_vector_x + math.cos(theta) * baseline_offset
            center_x = anchor_x - rotated_x
            center_y = baseline_y - rotated_y
            left = center_x - width_emu / 2
            top = center_y - height_emu / 2
            shape = self.slide.shapes.add_textbox(
                round(left), round(top), width_emu, height_emu
            )
        tf = shape.text_frame
        tf.clear()
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = tf.paragraphs[0]
        paragraph.alignment = alignment
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        run = paragraph.add_run()
        run.text = text
        run.font.name = self.font_name
        run.font.size = Pt(font_pt)
        run.font.italic = style.get("font-style", "").lower() in {"italic", "oblique"}
        weight = style.get("font-weight", "").lower()
        run.font.bold = weight in {"bold", "bolder", "600", "700", "800", "900"}
        text_color = color_value(style.get("fill", "#000000")) or RGBColor(0, 0, 0)
        run.font.color.rgb = text_color
        if vertical_text:
            tf._txBody.bodyPr.set("vert", "vert270" if angle < 0 else "vert")
        elif abs(angle) > 0.01:
            shape.rotation = angle % 360
        self.created.append(shape)

    def walk(
        self,
        element,
        parent_matrix: Matrix = IDENTITY,
        inherited_style: dict[str, str] | None = None,
        override_style: dict[str, str] | None = None,
        from_use: bool = False,
        inherited_clip: tuple[float, float, float, float] | None = None,
    ) -> None:
        tag = local_name(element)
        if self.suppress_legends and tag == "g" and (element.get("id") or "").startswith("legend_"):
            return
        if tag in {"metadata", "style", "clipPath"}:
            return
        if tag == "defs" and not from_use:
            return
        style = dict(inherited_style or {})
        style.update(parse_style(element))
        if override_style:
            style.update(override_style)
        matrix = multiply(parent_matrix, parse_transform(element.get("transform")))
        if tag == "g" and element.get("id") == "axes_1" and self.axis_normalization_matrix:
            matrix = multiply(matrix, self.axis_normalization_matrix)
        active_clip = inherited_clip
        clip_value = element.get("clip-path") or ""
        clip_match = re.search(r"url\(\s*#([^\s)]+)\s*\)", clip_value)
        if clip_match:
            raw_clip = self.clip_rects.get(clip_match.group(1))
            if raw_clip:
                corners = [
                    transform_point(matrix, raw_clip[0], raw_clip[1]),
                    transform_point(matrix, raw_clip[2], raw_clip[1]),
                    transform_point(matrix, raw_clip[2], raw_clip[3]),
                    transform_point(matrix, raw_clip[0], raw_clip[3]),
                ]
                xs = [point[0] for point in corners]
                ys = [point[1] for point in corners]
                active_clip = (min(xs), min(ys), max(xs), max(ys))
        if tag == "use":
            href = element.get(f"{{{XLINK_NS}}}href") or element.get("href")
            if href and href.startswith("#") and href[1:] in self.ids:
                use_matrix = multiply(matrix, (1, 0, 0, 1, numeric(element.get("x")), numeric(element.get("y"))))
                self.walk(
                    self.ids[href[1:]],
                    use_matrix,
                    inherited_style,
                    style,
                    from_use=True,
                    inherited_clip=active_clip,
                )
            return
        if tag == "path" and element.get("d"):
            self.add_path(element.get("d"), matrix, style, active_clip)
            return
        if tag == "rect":
            self.add_rect(element, matrix, style)
            return
        if tag == "line":
            self.add_line(element, matrix, style)
            return
        if tag == "polyline":
            self.add_poly(element.get("points", ""), matrix, style, False)
            return
        if tag == "polygon":
            self.add_poly(element.get("points", ""), matrix, style, True)
            return
        if tag == "circle":
            self.add_ellipse(element, matrix, style, True)
            return
        if tag == "ellipse":
            self.add_ellipse(element, matrix, style, False)
            return
        if tag == "text":
            self.add_text(element, matrix, style)
            return
        if tag == "image":
            self.image_count += 1
            return
        for child in element:
            self.walk(
                child,
                matrix,
                style,
                override_style,
                from_use=from_use,
                inherited_clip=active_clip,
            )

    def render(self, group_name: str):
        self.add_panel_anchor()
        self.add_common_axis_anchor()
        self.walk(self.root)
        self.add_category_proxy_markers()
        if self.image_count:
            raise ValueError(
                f"{self.svg_path} contains {self.image_count} raster <image> elements; "
                "editable mode rejects rasterized SVG content"
            )
        if not self.created:
            raise ValueError(f"No editable elements found in {self.svg_path}")
        group = self.slide.shapes.add_group_shape(self.created)
        set_shape_name(group, group_name)
        return group


def delete_all_slides(prs: Presentation) -> None:
    slide_ids = prs.slides._sldIdLst
    for slide_id in list(slide_ids):
        r_id = slide_id.rId
        slide_ids.remove(slide_id)
        prs.part.drop_rel(r_id)


def choose_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if layout.name == "1_\u7a7a\u767d":
            return layout
    return prs.slide_layouts[-1]


def format_title(slide, title: str) -> None:
    placeholders = [p for p in slide.placeholders if p.has_text_frame]
    if placeholders:
        shape = placeholders[0]
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), 0, Inches(52), Inches(2.9809))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string("785F8E")
        shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(120)
    r.font.color.rgb = RGBColor(255, 255, 255)


def add_boundary(slide) -> None:
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), 0, Inches(52), Inches(52))
    border.fill.background()
    border.line.color.rgb = RGBColor.from_string("C00000")
    border.line.width = Pt(10)
    set_shape_name(border, "Delivery boundary 2-54in")


def read_legend_config(path: Path | None) -> dict | None:
    if path is None:
        return None
    data = json.loads(path.read_text(encoding="utf-8"))
    items = data.get("items")
    if not isinstance(items, list) or not items:
        raise ValueError(f"Legend config {path} must contain a non-empty items list")
    for item in items:
        if not isinstance(item, dict) or not item.get("label") or not item.get("color"):
            raise ValueError(f"Every legend item in {path} needs label and color")
    return data


def add_shared_legend(
    slide,
    config: dict,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str,
) -> int:
    items = config["items"]
    font_pt = float(config.get("font_size", 72.0))
    marker_size = float(config.get("marker_size", 0.42))
    marker_gap = float(config.get("marker_gap", 0.22))
    item_gap = float(config.get("item_gap", 0.80))
    # Give native 72 pt labels enough intrinsic width for both PowerPoint and
    # stricter fallback renderers. Disabling wrap alone is not honored by all
    # office engines when the box is close to the measured glyph width.
    label_widths = [max(2.4, len(str(item["label"])) * font_pt * 0.68 / 72.0) for item in items]
    item_widths = [marker_size + marker_gap + label_width for label_width in label_widths]
    total_width = sum(item_widths) + item_gap * max(0, len(items) - 1)
    cursor = left + (width - total_width) / 2
    created = 0
    text_height = max(0.9, font_pt * 1.25 / 72.0)
    marker_top = top + (height - marker_size) / 2
    text_top = top + (height - text_height) / 2
    for index, (item, item_width, label_width) in enumerate(zip(items, item_widths, label_widths), start=1):
        marker_kind = str(item.get("marker", "circle")).lower()
        item_color = color_value(str(item["color"])) or RGBColor(0, 0, 0)
        if marker_kind == "line":
            marker_y = top + height / 2
            marker = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(cursor),
                Inches(marker_y),
                Inches(cursor + marker_size),
                Inches(marker_y),
            )
            marker.line.color.rgb = item_color
            marker.line.width = Pt(float(item.get("line_width", 5.0)))
        else:
            marker_shape = MSO_SHAPE.RECTANGLE if marker_kind in {"square", "rect"} else MSO_SHAPE.OVAL
            marker = slide.shapes.add_shape(
                marker_shape,
                Inches(cursor),
                Inches(marker_top),
                Inches(marker_size),
                Inches(marker_size),
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = item_color
            marker.line.color.rgb = RGBColor.from_string(str(item.get("line_color", "333333")).lstrip("#"))
            marker.line.width = Pt(float(item.get("line_width", 1.5)))
        set_shape_name(marker, f"Shared legend marker {index}: {item['label']}")
        text_left = cursor + marker_size + marker_gap
        textbox = slide.shapes.add_textbox(
            Inches(text_left),
            Inches(text_top),
            Inches(label_width),
            Inches(text_height),
        )
        textbox.text_frame.clear()
        textbox.text_frame.margin_left = 0
        textbox.text_frame.margin_right = 0
        textbox.text_frame.margin_top = 0
        textbox.text_frame.margin_bottom = 0
        textbox.text_frame.word_wrap = False
        textbox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = textbox.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = str(item["label"])
        run.font.name = font_name
        run.font.size = Pt(font_pt)
        run.font.color.rgb = RGBColor.from_string(str(config.get("text_color", "000000")).lstrip("#"))
        set_shape_name(textbox, f"Shared legend label {index}: {item['label']}")
        cursor += item_width + item_gap
        created += 2
    return created


def add_vertical_shared_legend(
    slide,
    config: dict,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str,
) -> int:
    """Build a native vertical legend inside a reserved grid cell."""
    items = config["items"]
    font_pt = float(config.get("font_size", 64.0))
    marker_slot = float(config.get("marker_size", 1.15))
    marker_gap = float(config.get("marker_gap", 0.30))
    padding_x = float(config.get("padding_x", 0.35))
    padding_y = float(config.get("padding_y", 0.35))
    usable_height = max(0.1, height - 2 * padding_y)
    configured_item_height = config.get("item_height")
    if configured_item_height is not None:
        item_height = min(float(configured_item_height), usable_height / len(items))
        legend_top = top + (height - item_height * len(items)) / 2
    else:
        item_height = usable_height / len(items)
        legend_top = top + padding_y
    text_height = min(item_height, max(0.75, font_pt * 1.25 / 72.0))
    marker_left = left + padding_x
    text_left = marker_left + marker_slot + marker_gap
    text_width = max(0.5, left + width - padding_x - text_left)
    created = 0

    for index, item in enumerate(items, start=1):
        row_top = legend_top + (index - 1) * item_height
        row_center = row_top + item_height / 2
        marker_kind = str(item.get("marker", "circle")).lower()
        item_color = color_value(str(item["color"])) or RGBColor(0, 0, 0)
        line_width = float(item.get("line_width", 5.0))

        if marker_kind == "line":
            marker = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(marker_left),
                Inches(row_center),
                Inches(marker_left + marker_slot),
                Inches(row_center),
            )
            marker.line.color.rgb = item_color
            marker.line.width = Pt(line_width)
            if str(item.get("dash", "")).lower() in {"dash", "dashed"}:
                marker.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        else:
            marker_size = min(
                float(item.get("shape_size", marker_slot * 0.62)),
                item_height * 0.62,
            )
            marker_shape = {
                "square": MSO_SHAPE.RECTANGLE,
                "rect": MSO_SHAPE.RECTANGLE,
                "diamond": MSO_SHAPE.DIAMOND,
            }.get(marker_kind, MSO_SHAPE.OVAL)
            marker = slide.shapes.add_shape(
                marker_shape,
                Inches(marker_left + (marker_slot - marker_size) / 2),
                Inches(row_center - marker_size / 2),
                Inches(marker_size),
                Inches(marker_size),
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = item_color
            marker.line.color.rgb = RGBColor.from_string(
                str(item.get("line_color", "333333")).lstrip("#")
            )
            marker.line.width = Pt(float(item.get("line_width", 1.5)))
        set_shape_name(marker, f"Shared legend marker {index}: {item['label']}")

        textbox = slide.shapes.add_textbox(
            Inches(text_left),
            Inches(row_center - text_height / 2),
            Inches(text_width),
            Inches(text_height),
        )
        textbox.text_frame.clear()
        textbox.text_frame.margin_left = 0
        textbox.text_frame.margin_right = 0
        textbox.text_frame.margin_top = 0
        textbox.text_frame.margin_bottom = 0
        textbox.text_frame.word_wrap = False
        textbox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = textbox.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = str(item["label"])
        run.font.name = font_name
        run.font.size = Pt(font_pt)
        run.font.color.rgb = RGBColor.from_string(
            str(config.get("text_color", "000000")).lstrip("#")
        )
        set_shape_name(textbox, f"Shared legend label {index}: {item['label']}")
        created += 2
    return created


def build_deck(args) -> tuple[int, int]:
    prs = Presentation(str(args.template))
    if prs.slide_width != 56 * EMU_PER_INCH or prs.slide_height != 56 * EMU_PER_INCH:
        raise ValueError(
            f"Template canvas is {prs.slide_width / EMU_PER_INCH:.3f} x "
            f"{prs.slide_height / EMU_PER_INCH:.3f} in; expected 56 x 56 in"
        )
    delete_all_slides(prs)
    slide = prs.slides.add_slide(choose_layout(prs))
    format_title(slide, args.title)
    svg_paths = sorted(args.svg_dir.glob("panel-*.svg"))
    if not 2 <= len(svg_paths) <= 12:
        raise ValueError(f"Expected 2–12 SVG panels, found {len(svg_paths)} in {args.svg_dir}")
    content_left = 2.0
    content_top = 3.35
    content_width = 52.0
    content_height = 48.15
    geometries = [read_panel_geometry(path) for path in svg_paths]
    resolved_row_counts = (
        list(args.row_counts)
        if args.row_counts
        else balanced_row_counts(len(geometries), args.max_columns)
    )
    for row in args.suppress_shared_category_rows or []:
        if row < 1 or row > len(resolved_row_counts):
            raise ValueError(f"Shared-category suppression row is out of range: {row}")
    for panel in args.shared_category_proxy_panels or []:
        if panel < 1 or panel > len(svg_paths):
            raise ValueError(f"Shared-category proxy panel is out of range: {panel}")
    if args.last_row_bottom_margin is not None and args.last_row_bottom_margin < 0:
        raise ValueError("--last-row-bottom-margin cannot be negative")
    legend_config = read_legend_config(args.legend_config)
    if (args.suppress_shared_category_rows or args.shared_category_proxy_panels) and not legend_config:
        raise ValueError("Shared-category suppression/proxies require --legend-config")
    legend_position = (
        str(legend_config.get("position", "after-first-row")).lower()
        if legend_config
        else ""
    )
    legend_band_height = (
        args.legend_band_height
        if legend_config and legend_position == "after-first-row"
        else 0.0
    )
    boxes = adaptive_layout(
        geometries,
        left=content_left,
        top=content_top,
        width=content_width,
        height=content_height,
        max_columns=args.max_columns,
        min_gap_x=args.min_gap_x,
        min_gap_y=args.min_gap_y,
        legend_band_height=legend_band_height,
        lock_common_panel_size=args.lock_common_panel_size,
        stretch_row_panels=args.stretch_row_panels,
        stretch_rows=set(args.stretch_rows or []),
        stretch_row_weights=args.stretch_row_weights,
        row_counts=args.row_counts,
    )
    if args.fit_rows_to_width:
        selected_rows = set(args.fit_rows_to_width)
        invalid_rows = sorted(
            row for row in selected_rows if row < 1 or row > len(resolved_row_counts)
        )
        if invalid_rows:
            raise ValueError(
                "--fit-rows-to-width contains invalid one-based row indices: "
                + ", ".join(map(str, invalid_rows))
            )
        cursor = 0
        for row_index, row_count in enumerate(resolved_row_counts, start=1):
            if row_index in selected_rows:
                row_geometries = geometries[cursor : cursor + row_count]
                usable_width = content_width - args.min_gap_x * max(0, row_count - 1)
                row_height = usable_width / sum(
                    geometry.aspect for geometry in row_geometries
                )
                row_top = min(
                    top for _, top, _, _ in boxes[cursor : cursor + row_count]
                )
                row_boxes = []
                left = content_left
                for geometry in row_geometries:
                    panel_width = geometry.aspect * row_height
                    row_boxes.append((left, row_top, panel_width, row_height))
                    left += panel_width + args.min_gap_x
                boxes[cursor : cursor + row_count] = row_boxes
            cursor += row_count
    if args.pack_locked_rows or args.pack_rows:
        if args.pack_locked_rows and not args.lock_common_panel_size:
            raise ValueError("--pack-locked-rows requires --lock-common-panel-size")
        cursor = resolved_row_counts[0]
        previous_start = 0
        previous_count = resolved_row_counts[0]
        for row_count in resolved_row_counts[1:]:
            previous_bottom = max(
                top + height
                for _, top, _, height in boxes[previous_start : previous_start + previous_count]
            )
            current_top = min(top for _, top, _, _ in boxes[cursor : cursor + row_count])
            shift = previous_bottom + args.min_gap_y - current_top
            boxes[cursor : cursor + row_count] = [
                (left, top + shift, width, height)
                for left, top, width, height in boxes[cursor : cursor + row_count]
            ]
            previous_start = cursor
            previous_count = row_count
            cursor += row_count
        if args.center_packed_rows:
            block_top = min(top for _, top, _, _ in boxes)
            block_bottom = max(top + height for _, top, _, height in boxes)
            target_midpoint = (content_top + 52.0) / 2.0
            shift = target_midpoint - (block_top + block_bottom) / 2.0
            boxes = [
                (left, top + shift, width, height)
                for left, top, width, height in boxes
            ]
        elif args.pack_rows:
            block_top = min(top for _, top, _, _ in boxes)
            shift = content_top - block_top
            boxes = [
                (left, top + shift, width, height)
                for left, top, width, height in boxes
            ]
    elif args.center_packed_rows:
        raise ValueError("--center-packed-rows requires --pack-rows or --pack-locked-rows")
    if args.compact_shared_legend:
        if not legend_config or legend_position != "after-first-row" or len(resolved_row_counts) < 2:
            raise ValueError(
                "--compact-shared-legend requires an after-first-row shared legend and at least two rows"
            )
        first_row_count = resolved_row_counts[0]
        next_row_top = min(top for _, top, _, _ in boxes[first_row_count:])
        first_row_bottom = max(top + height for _, top, _, height in boxes[:first_row_count])
        target_bottom = next_row_top - args.compact_row_gap
        shift = target_bottom - first_row_bottom
        boxes[:first_row_count] = [
            (left, top + shift, width, height)
            for left, top, width, height in boxes[:first_row_count]
        ]
    if args.last_row_bottom_margin is not None:
        final_row_count = resolved_row_counts[-1]
        final_row_start = len(boxes) - final_row_count
        target_bottom = 52.0 - args.last_row_bottom_margin
        current_bottom = max(
            top + height for _, top, _, height in boxes[final_row_start:]
        )
        shift = target_bottom - current_bottom
        boxes[final_row_start:] = [
            (left, top + shift, width, height)
            for left, top, width, height in boxes[final_row_start:]
        ]
    legend_box = None
    if legend_config and legend_position in {"bottom-right", "last-row-right"}:
        final_row_count = resolved_row_counts[-1]
        if resolved_row_counts[0] != args.max_columns or final_row_count >= args.max_columns:
            raise ValueError(
                "bottom-right legend placement requires a complete first row and an incomplete final row"
            )
        final_row_start = len(boxes) - final_row_count
        for offset in range(final_row_count):
            reference_left, _, reference_width, _ = boxes[offset]
            _, panel_top, _, panel_height = boxes[final_row_start + offset]
            boxes[final_row_start + offset] = (
                reference_left,
                panel_top,
                reference_width,
                panel_height,
            )
        reference_left, _, reference_width, _ = boxes[args.max_columns - 1]
        _, final_top, _, final_height = boxes[-1]
        legend_box = (reference_left, final_top, reference_width, final_height)
    common_aspect = Counter(round(geometry.aspect, 6) for geometry in geometries).most_common(1)[0][0]
    titles = [read_panel_title(path) for path in svg_paths]
    panel_rows: list[int] = []
    for row_index, count in enumerate(resolved_row_counts, start=1):
        panel_rows.extend([row_index] * count)
    suppressed_category_rows = set(args.suppress_shared_category_rows or [])
    proxy_panels = set(args.shared_category_proxy_panels or [])
    shared_items = list(legend_config.get("items", [])) if legend_config else []
    ratio_panels = set(args.preserve_source_font_ratio_panels or [])
    invalid_ratio_panels = sorted(
        panel for panel in ratio_panels if panel < 1 or panel > len(svg_paths)
    )
    if invalid_ratio_panels:
        raise ValueError(
            "--preserve-source-font-ratio-panels contains out-of-range panel "
            f"indices: {invalid_ratio_panels}"
        )
    panel_detail_sizes: dict[int, float] = {}
    raw_panel_detail_sizes = args.panel_detail_font_sizes or []
    if len(raw_panel_detail_sizes) % 2:
        raise ValueError(
            "--panel-detail-font-sizes requires PANEL PT pairs, for example "
            "--panel-detail-font-sizes 2 44 8 36"
        )
    for panel_value, size_value in zip(
        raw_panel_detail_sizes[::2], raw_panel_detail_sizes[1::2]
    ):
        panel = int(panel_value)
        if panel_value != panel or panel < 1 or panel > len(svg_paths):
            raise ValueError(
                "--panel-detail-font-sizes contains an invalid one-based panel "
                f"index: {panel_value}"
            )
        if size_value <= 0:
            raise ValueError(
                "--panel-detail-font-sizes requires positive point sizes"
            )
        panel_detail_sizes[panel] = size_value
    title_font_scales: list[float] = []
    cursor = 0
    for count in resolved_row_counts:
        row_scales = []
        for title, box in zip(titles[cursor : cursor + count], boxes[cursor : cursor + count]):
            available_width_pt = max(1.0, (box[2] - 0.55) * 72.0)
            estimated_width_pt = args.font_size * (len(title) * 0.57 + 0.25)
            row_scales.append(min(1.0, available_width_pt / estimated_width_pt))
        common_row_scale = (
            1.0
            if args.font_levels or args.preserved_ratio_role_levels
            else max(0.72, min(row_scales, default=1.0))
        )
        title_font_scales.extend([common_row_scale] * count)
        cursor += count
    groups = []
    for index, (svg_path, geometry, box, title_font_scale) in enumerate(
        zip(svg_paths, geometries, boxes, title_font_scales)
    ):
        preserve_panel_font_ratios = (
            args.preserve_source_font_ratios or index + 1 in ratio_panels
        )
        panel_font_levels = tuple(args.font_levels) if args.font_levels else None
        if panel_font_levels and index + 1 in panel_detail_sizes:
            panel_font_levels = (
                panel_font_levels[0],
                panel_detail_sizes[index + 1],
            )
        normalize_axis_box = None
        if args.normalize_common_axes and abs(geometry.aspect - common_aspect) <= args.common_aspect_tolerance:
            normalize_axis_box = tuple(args.common_axis_box)
        renderer = SvgPanelRenderer(
            slide,
            svg_path,
            box,
            target_font_pt=args.font_size,
            font_name=args.font,
            normalize_titles=not args.no_normalize_titles,
            preserve_source_font_ratios=preserve_panel_font_ratios,
            preserved_ratio_font_levels=(
                tuple(args.preserved_ratio_font_levels)
                if args.preserved_ratio_font_levels
                else None
            ),
            preserved_ratio_role_levels=(
                tuple(args.preserved_ratio_role_levels)
                if args.preserved_ratio_role_levels
                else None
            ),
            normalize_axis_box=normalize_axis_box,
            suppress_legends=bool(legend_config),
            compact_rotated_labels=bool(args.normalize_common_axes and normalize_axis_box is None),
            title_font_scale=title_font_scale,
            font_levels=(
                None
                if preserve_panel_font_ratios
                else panel_font_levels
            ),
            dense_font_pt=args.dense_font_size,
            very_dense_font_pt=args.very_dense_font_size,
            min_stroke_width_pt=args.min_stroke_width,
            axis_stroke_width_pt=args.axis_stroke_width,
            suppressed_texts={
                alias
                for item in shared_items
                for alias in [str(item["label"]), *map(str, item.get("aliases", []))]
            }
            if panel_rows[index] in suppressed_category_rows
            else set(),
            category_proxy_items=shared_items if index + 1 in proxy_panels else [],
            category_proxy_gap_in=args.category_proxy_gap,
        )
        groups.append(renderer.render(f"Editable SVG panel {index + 1:02d}: {svg_path.stem}"))
    legend_elements = 0
    if legend_config and legend_position in {"bottom-right", "last-row-right"}:
        assert legend_box is not None
        legend_elements = add_vertical_shared_legend(
            slide,
            legend_config,
            left=legend_box[0],
            top=legend_box[1],
            width=legend_box[2],
            height=legend_box[3],
            font_name=args.font,
        )
    elif legend_config:
        first_row_count = resolved_row_counts[0]
        first_row_bottom = max(top + height for _, top, _, height in boxes[:first_row_count])
        legend_top = (
            first_row_bottom - args.compact_legend_overlap
            if args.compact_shared_legend
            else first_row_bottom + args.min_gap_y * 0.25
        )
        legend_elements = add_shared_legend(
            slide,
            legend_config,
            left=content_left,
            top=legend_top,
            width=content_width,
            height=legend_band_height,
            font_name=args.font,
        )
    add_boundary(slide)
    if args.notes_file:
        slide.notes_slide.notes_text_frame.text = args.notes_file.read_text(encoding="utf-8")
    args.output.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(prefix="editable-svg-", suffix=".pptx", dir="/private/tmp", delete=False) as tmp:
        temp_path = Path(tmp.name)
    prs.save(temp_path)
    shutil.copy2(temp_path, args.output)
    temp_path.unlink(missing_ok=True)
    return len(svg_paths), sum(len(group.shapes) for group in groups) + legend_elements


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--template", type=Path, required=True)
    parser.add_argument("--svg-dir", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--title", default="Fig S2")
    parser.add_argument("--font", default="Calibri")
    parser.add_argument("--font-size", type=float, default=88.0)
    parser.add_argument(
        "--font-levels",
        type=float,
        nargs=2,
        metavar=("HEADING_PT", "PLOT_TEXT_PT"),
        help="Use HEADING_PT for panel titles and one uniform PLOT_TEXT_PT for every other chart label; with the 120 pt slide title this yields exactly three font levels",
    )
    parser.add_argument(
        "--panel-detail-font-sizes",
        type=float,
        nargs="+",
        metavar=("PANEL", "PT"),
        help="Override the non-heading font size for selected one-based panels using PANEL PT pairs",
    )
    parser.add_argument("--max-columns", type=int, choices=(2, 3, 4), default=3)
    parser.add_argument(
        "--row-counts",
        type=int,
        nargs="+",
        help="Explicit panels per row, for example --row-counts 4 2 3; values must sum to the SVG count",
    )
    parser.add_argument("--min-gap-x", type=float, default=0.80)
    parser.add_argument("--min-gap-y", type=float, default=0.80)
    parser.add_argument(
        "--pack-locked-rows",
        action="store_true",
        help="Pack locked-size rows using --min-gap-y as the actual panel-box gap",
    )
    parser.add_argument(
        "--pack-rows",
        action="store_true",
        help="Pack natural-aspect rows using --min-gap-y as the actual panel-box gap",
    )
    parser.add_argument(
        "--center-packed-rows",
        action="store_true",
        help="Vertically center a packed row block between the content top and the 52-inch delivery boundary",
    )
    parser.add_argument(
        "--compact-shared-legend",
        action="store_true",
        help="Overlay the shared legend into unused lower space in the first-row panels and move that row next to row two",
    )
    parser.add_argument(
        "--compact-row-gap",
        type=float,
        default=0.45,
        help="Vertical gap in inches between the compacted first-row panel boxes and row two",
    )
    parser.add_argument(
        "--compact-legend-overlap",
        type=float,
        default=1.115,
        help="How far the shared-legend band overlaps the bottom of the compacted first-row panel boxes, in inches",
    )
    parser.add_argument(
        "--last-row-bottom-margin",
        type=float,
        help="Align the final row to this margin above the 52-inch delivery boundary",
    )
    parser.add_argument("--lock-common-panel-size", action="store_true")
    parser.add_argument(
        "--stretch-row-panels",
        action="store_true",
        help="Stretch panels horizontally into equal-width row cells so every populated row fills the content width",
    )
    parser.add_argument(
        "--stretch-rows",
        type=int,
        nargs="+",
        metavar="ROW",
        help="One-based row numbers to fill horizontally while preserving relative source-width proportions",
    )
    parser.add_argument(
        "--fit-rows-to-width",
        type=int,
        nargs="+",
        metavar="ROW",
        help="Fit selected one-based rows to the full content width while preserving every source panel aspect ratio; row height grows as needed",
    )
    parser.add_argument(
        "--stretch-row-weights",
        type=float,
        nargs="+",
        metavar="WEIGHT",
        help="Optional per-panel multipliers for a selected stretched row, for example 1 1.25 1",
    )
    parser.add_argument("--normalize-common-axes", action="store_true")
    parser.add_argument(
        "--common-axis-box",
        type=float,
        nargs=4,
        metavar=("LEFT", "RIGHT", "TOP", "BOTTOM"),
        default=(0.32, 0.98, 0.14, 0.74),
        help="Normalized SVG fractions for the shared non-exception axes rectangle",
    )
    parser.add_argument("--common-aspect-tolerance", type=float, default=0.01)
    parser.add_argument(
        "--dense-font-size",
        type=float,
        help="Use this smaller non-heading font size only for multi-axes panels such as heatmaps with color bars",
    )
    parser.add_argument(
        "--very-dense-font-size",
        type=float,
        help="Use this non-heading font size for panels with at least 35 SVG text elements",
    )
    parser.add_argument(
        "--min-stroke-width",
        type=float,
        help="Minimum width in points for non-black vector strokes",
    )
    parser.add_argument(
        "--axis-stroke-width",
        type=float,
        help="Fixed width in points for black axis and tick strokes",
    )
    parser.add_argument(
        "--suppress-shared-category-rows",
        type=int,
        nargs="+",
        metavar="ROW",
        help="Remove category tick labels that duplicate shared-legend labels in the selected one-based rows",
    )
    parser.add_argument(
        "--shared-category-proxy-panels",
        type=int,
        nargs="+",
        metavar="PANEL",
        help="Add color-only column markers above the primary axes in selected one-based panels",
    )
    parser.add_argument(
        "--category-proxy-gap",
        type=float,
        default=0.60,
        help="Gap in inches between proxy-marker bottoms and the primary axes top",
    )
    parser.add_argument("--legend-config", type=Path)
    parser.add_argument("--legend-band-height", type=float, default=1.60)
    parser.add_argument(
        "--notes-file",
        type=Path,
        help="Optional UTF-8 text file stored in the generated slide notes, useful for exact values omitted from a presentation-first visual",
    )
    parser.add_argument("--no-normalize-titles", action="store_true")
    parser.add_argument(
        "--preserve-source-font-ratios",
        action="store_true",
        help="Keep explicitly smaller SVG ticks/legends proportional while primary text uses --font-size",
    )
    parser.add_argument(
        "--preserve-source-font-ratio-panels",
        type=int,
        nargs="+",
        metavar="PANEL",
        help="Use proportional source font sizes only for selected one-based panels while --font-levels remains active elsewhere",
    )
    parser.add_argument(
        "--preserved-ratio-font-levels",
        type=float,
        nargs=3,
        metavar=("PRIMARY_PT", "SECONDARY_PT", "MICRO_PT"),
        help="Quantize proportional-font panels to three non-heading point sizes instead of retaining every source size",
    )
    parser.add_argument(
        "--preserved-ratio-role-levels",
        type=float,
        nargs=4,
        metavar=("SECTION_PT", "AXIS_PT", "BODY_PT", "MICRO_PT"),
        help="Map standardized source-size roles to four deliberate non-heading PowerPoint sizes; intended for SVGs using 40/34/30/24/18-20 pt heading/section/axis/body/micro source text",
    )
    args = parser.parse_args()
    if args.preserved_ratio_font_levels and args.preserved_ratio_role_levels:
        parser.error(
            "--preserved-ratio-font-levels and --preserved-ratio-role-levels "
            "are mutually exclusive"
        )
    panels, elements = build_deck(args)
    print(f"Created {args.output} with {panels} editable SVG groups and {elements} native elements")


if __name__ == "__main__":
    main()
